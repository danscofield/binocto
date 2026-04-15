"""Two-slide context deck: the state of the art before agentic AI."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (matches main deck) ───────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x0F, 0x1A)
ACCENT  = RGBColor(0x5B, 0x8D, 0xFF)
ACCENT2 = RGBColor(0x7C, 0xE8, 0xB4)
WARN    = RGBColor(0xFF, 0x7C, 0x5C)
GOLD    = RGBColor(0xFF, 0xCC, 0x55)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xCC, 0xD6, 0xF1)
DIM     = RGBColor(0x77, 0x88, 0xAA)
BOX_BG  = RGBColor(0x1A, 0x1E, 0x35)
BOX_BDR = RGBColor(0x2E, 0x3A, 0x5C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Primitives ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = BG

def rect(slide, l, t, w, h, fill=BOX_BG, border=BOX_BDR, bpt=1.0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(bpt)
    return s

def bar(slide, l, t, w, h, color=ACCENT):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def title_block(slide, title, subtitle=None):
    bg(slide)
    bar(slide, 0.5, 0.58, 12.33, 0.06)
    txt(slide, title,    0.5, 0.1,  12.33, 0.55, size=32, bold=True)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.68, 12.33, 0.38, size=16, color=LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — The old world
# ═══════════════════════════════════════════════════════════════════════════════

def slide_old_world(prs):
    sl = blank(prs)
    title_block(sl,
        "Vulnerability Research Before Agentic AI",
        "Expert-bottlenecked tools that don't compose — and don't scale")

    # ── Left column: traditional approaches ──────────────────────────────────
    txt(sl, "Traditional approaches", 0.5, 1.12, 5.9, 0.38,
        size=15, bold=True, color=ACCENT)

    approaches = [
        ("Manual reverse engineering",
         "Ghidra / IDA Pro + human expert.\nDays to weeks per binary. Doesn't scale."),
        ("Rule-based static analysis",
         "Coverity, CodeQL, Semgrep.\nPatterns written by humans — miss novel code paths,\nhigh false-positive rates, can't reason about context."),
        ("Symbolic execution",
         "KLEE, angr, Triton.\nModels all possible paths algebraically.\nPath explosion: state space doubles at every branch.\nPractical ceiling ~10k lines of code."),
        ("Fuzzing",
         "AFL++, LibFuzzer, honggfuzz.\nExcellent at crashing — poor at understanding\nwhat configuration is required to reach a path.\nBlind to semantic preconditions."),
    ]
    cy = 1.52
    cols = [ACCENT, ACCENT2, WARN, GOLD]
    for i, (name, body) in enumerate(approaches):
        bx = rect(sl, 0.5, cy, 5.9, 1.2)
        bar(sl, 0.5, cy, 0.16, 1.2, color=cols[i])
        txt(sl, name, 0.8, cy+0.08, 5.45, 0.38, size=13, bold=True, color=cols[i])
        txt(sl, body, 0.8, cy+0.46, 5.45, 0.65, size=12, color=LIGHT)
        cy += 1.28

    # ── Right column: fundamental limits ─────────────────────────────────────
    txt(sl, "Why they fall short", 6.65, 1.12, 6.2, 0.38,
        size=15, bold=True, color=WARN)

    limits = [
        (WARN,   "No semantic reasoning",
                 "Static tools match patterns; they cannot ask: if this flag were absent, would the "
                 "attacker still reach the sink? That question requires understanding intent, "
                 "not just code structure."),
        (WARN,   "Doesn't compose across abstraction layers",
                 "A vulnerability may require: a CLI flag → sets a global → tested in a dispatcher → "
                 "calls a handler → passes user input to system(). Each tool sees one layer. "
                 "None connect them automatically."),
        (WARN,   "Configuration blindness",
                 "Symbolic execution explores code paths but has no model of what argv, config files, "
                 "or prior HTTP requests mean. Preconditions — the thing defenders actually need — "
                 "are invisible to the tool."),
        (DIM,    "Expert dependency never eliminated",
                 "Every tool in this stack produces raw output that requires a skilled analyst to "
                 "interpret, triage, and connect into an actionable finding. "
                 "Automation assists; it does not replace."),
    ]
    cy2 = 1.52
    for col, name, body in limits:
        bx = rect(sl, 6.65, cy2, 6.2, 1.2)
        bar(sl, 6.65, cy2, 0.16, 1.2, color=col)
        txt(sl, name, 6.95, cy2+0.08, 5.75, 0.38, size=13, bold=True, color=col)
        txt(sl, body, 6.95, cy2+0.46, 5.75, 0.65, size=12, color=LIGHT)
        cy2 += 1.28

    # ── Divider ──────────────────────────────────────────────────────────────
    bar(sl, 6.35, 1.1, 0.04, 5.85, color=BOX_BDR)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — DARPA
# ═══════════════════════════════════════════════════════════════════════════════

def slide_darpa(prs):
    sl = blank(prs)
    title_block(sl,
        "DARPA's Decade-Long Push  —  Still an Open Problem",
        "Hundreds of millions of dollars, world-class teams, eight years of competition")

    # ── Timeline boxes ───────────────────────────────────────────────────────
    programs = [
        ("2014 – 2016",
         "CGC\nCyber Grand Challenge",
         ACCENT,
         "First fully automated capture-the-flag tournament. Seven finalist systems competed "
         "on a custom sandboxed OS (DECREE) using purpose-built, toy-scale challenge binaries "
         "— not real-world software. Mayhem (CMU ForAllSecure) won.\n\n"
         "Proved automation was possible in a controlled box. "
         "Did not generalise beyond the competition environment."),
        ("2016 – 2023",
         "HACCS & related\nHarnessing Autonomy for\nCybersecurity",
         GOLD,
         "Follow-on DARPA programs explored scaling automated analysis to larger and more "
         "realistic targets — networked systems, multi-component software, real OS kernels. "
         "Research advanced the state of the art in binary analysis, taint tracking, and "
         "automated exploit generation.\n\n"
         "Scaling beyond lab conditions remained the central unsolved challenge."),
        ("2023 – present",
         "AIxCC\nAI Cyber Challenge",
         ACCENT2,
         "Two-year competition with ~$18.5 M in prizes, backed by Anthropic, Google, Microsoft, "
         "and OpenZeppelin. Targets real open-source software: Linux kernel, nginx, SQLite.\n\n"
         "Marks the transition to LLM-assisted analysis. Semifinal results (2024) showed "
         "meaningful progress — but top systems still required significant human scaffolding "
         "and struggled with novel vulnerability classes at arbitrary scale."),
    ]

    bw = 3.85
    for i, (period, name, col, body) in enumerate(programs):
        x = 0.45 + i*(bw + 0.19)
        bx = rect(sl, x, 1.22, bw, 5.6)
        # top colour strip
        bar(sl, x, 1.22, bw, 0.28, color=col)
        txt(sl, period, x+0.12, 1.25, bw-0.2, 0.28, size=11, bold=False, color=BG)
        txt(sl, name,   x+0.12, 1.55, bw-0.2, 0.72, size=14, bold=True, color=col)
        txt(sl, body,   x+0.12, 2.32, bw-0.2, 4.2,  size=12, color=LIGHT)

    # ── Timeline connector ───────────────────────────────────────────────────
    bar(sl, 0.45, 1.35, 12.33, 0.04, color=DIM)
    for i in range(3):
        cx = 0.45 + i*(bw+0.19) + bw/2 - 0.12
        bar(sl, cx, 1.22, 0.25, 0.25, color=programs[i][2])

    # ── Bottom callout ───────────────────────────────────────────────────────
    bx = rect(sl, 0.45, 6.88, 12.4, 0.52, fill=RGBColor(0x14,0x08,0x08), border=WARN, bpt=1.2)
    txt(sl,
        "Eight years.  Hundreds of millions of dollars.  The world's best security researchers and AI labs.  "
        "Automated vulnerability analysis at general scale — across arbitrary real-world binaries, without human guidance — "
        "remains an open problem.  Agentic LLMs are the first technique to show genuine compositional reasoning across "
        "the full analysis chain.",
        0.6, 6.92, 12.1, 0.44, size=12.5, color=LIGHT, italic=False)


# ═══════════════════════════════════════════════════════════════════════════════
# Build
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_old_world(prs)
    slide_darpa(prs)
    out = "/home/ubuntu/anthropic/context_before_agentic_ai.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
