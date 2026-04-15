"""Generate a PowerPoint deck summarising the binary vulnerability analysis agent project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x0F, 0x1A)   # near-black blue
ACCENT    = RGBColor(0x5B, 0x8D, 0xFF)   # bright blue
ACCENT2   = RGBColor(0x7C, 0xE8, 0xB4)   # mint green
WARN      = RGBColor(0xFF, 0x7C, 0x5C)   # coral
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xD6, 0xF1)   # pale blue-grey
DIM       = RGBColor(0x77, 0x88, 0xAA)
BOX_BG    = RGBColor(0x1A, 0x1E, 0x35)   # slightly lighter than BG
BOX_BDR   = RGBColor(0x2E, 0x3A, 0x5C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def bg(slide):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def box(slide, l, t, w, h, fill=BOX_BG, border=BOX_BDR, border_pt=1.0):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(border_pt)
    return shape


def accent_bar(slide, t=0.58, h=0.06):
    """Thin horizontal accent bar below the title."""
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(t), Inches(12.33), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def slide_title(slide, title, subtitle=None):
    bg(slide)
    accent_bar(slide)
    txt(slide, title, 0.5, 0.12, 12.33, 0.55,
        size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.72, 12.33, 0.4,
            size=16, color=LIGHT, align=PP_ALIGN.LEFT)


def bullet_block(slide, items, l, t, w, col=LIGHT, size=17, gap=0.36, leader="›  "):
    """Render a list of bullet strings, each on its own text box."""
    cy = t
    for item in items:
        indent = 0
        text = item
        if item.startswith("  "):
            indent = 0.35
            text = item.lstrip()
            sz = size - 2
            c  = DIM
        else:
            sz = size
            c  = col
        txt(slide, leader + text, l + indent, cy, w - indent, gap,
            size=sz, color=c)
        cy += gap
    return cy


def colored_box_txt(slide, label, value, l, t, w, h=1.0,
                    lcolor=ACCENT, vsize=36):
    b = box(slide, l, t, w, h)
    txt(slide, label, l+0.1, t+0.08, w-0.2, 0.35, size=12, color=DIM)
    txt(slide, value, l+0.1, t+0.35, w-0.2, h-0.4,
        size=vsize, bold=True, color=lcolor, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# Slides
# ═══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    sl = blank_slide(prs)
    bg(sl)
    # large decorative bar
    shape = sl.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(13.33), Inches(1.7))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x14, 0x18, 0x2E)
    shape.line.fill.background()

    txt(sl, "Binary Vulnerability Analysis Agent",
        0.7, 1.6, 12.0, 1.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(sl, "Automated static analysis of compiled HTTP servers\nfor exploitable command-injection preconditions",
        0.7, 2.95, 11.5, 1.0, size=22, color=LIGHT, align=PP_ALIGN.LEFT)

    accent_bar(sl, t=2.6, h=0.07)

    txt(sl, "Agentic LLM  ·  Joern CPG  ·  LangGraph  ·  50-sample corpus evaluation",
        0.7, 6.05, 12.0, 0.55, size=14, color=DIM, align=PP_ALIGN.LEFT)


def slide_problem(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Problem", "What are we trying to do?")

    bullet_block(sl, [
        "Given a compiled C binary (stripped HTTP server) — no source code",
        "Determine whether a command-injection vulnerability is exploitable",
        "Identify the exact configuration preconditions required to reach the sink",
        "",
        "Preconditions may be:",
        "  CLI flags that must be present at startup  (--exec-mode)",
        "  Config-file keys with specific values  (exec_mode = 1)",
        "  A prior HTTP request that initialises runtime state  (POST /exec/init)",
        "",
        "The corpus is 190+ synthetically patched binaries across 12 HTTP server codebases",
        "Ground truth: verified PoC shell scripts + JSON precondition lists",
    ], l=0.6, t=1.15, w=12.0, size=18, gap=0.38)


def slide_approach(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Approach", "Agentic static analysis — no binary execution")

    # Three principle boxes
    for i, (title, body) in enumerate([
        ("Static only",
         "Joern code property graph,\nobjdump disassembly, string extraction.\nNo dynamic execution or sandboxing."),
        ("LLM reasoning",
         "Claude traces call graphs, reads assembly,\nidentifies gate variables and their config\nsource across function boundaries."),
        ("Verified hypothesis",
         "A synthesise → verify loop runs up to 3×.\nEach iteration refines the precondition\nhypothesis against binary evidence."),
    ]):
        bx = box(sl, 0.45 + i*4.25, 1.35, 4.0, 2.2)
        txt(sl, title, 0.65 + i*4.25, 1.45, 3.7, 0.45, size=16, bold=True, color=ACCENT)
        txt(sl, body,  0.65 + i*4.25, 1.95, 3.7, 1.45, size=15, color=LIGHT)

    # Arrow pipeline
    steps = ["find_sinks", "check_input_reach", "trace_gate +\nclassify_source",
             "trace_input_path", "assess_sanitization", "synthesize\npreconditions",
             "causal_verify", "aggregate"]
    sw = 1.38
    gap = 0.06
    sy = 3.95
    sh = 0.9
    for j, s in enumerate(steps):
        x = 0.38 + j*(sw + gap)
        bx = box(sl, x, sy, sw, sh, fill=BOX_BG, border=ACCENT, border_pt=0.8)
        txt(sl, s, x+0.05, sy+0.1, sw-0.1, sh-0.15, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
        if j < len(steps)-1:
            txt(sl, "›", x+sw+gap/2-0.09, sy+0.28, 0.2, 0.35, size=16, color=ACCENT, bold=True)

    txt(sl, "LangGraph — parallel per-sink fan-out via Send API",
        0.45, 5.05, 12.3, 0.4, size=13, color=DIM, italic=True)

    txt(sl, "Synthesise / Verify loop (up to 3 retries)",
        7.55, 4.88, 4.2, 0.3, size=11, color=ACCENT2, italic=True)
    # bracket
    bracket = sl.shapes.add_shape(1, Inches(7.5), Inches(5.18), Inches(4.2), Inches(0.04))
    bracket.fill.solid(); bracket.fill.fore_color.rgb = ACCENT2
    bracket.line.fill.background()


def slide_causal(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Causal Link Model", "All four links must hold for an Exploitable verdict")

    links = [
        ("1  config → variable",
         "The CLI flag / config key / HTTP init request\nactually writes the gate variable into memory."),
        ("2  variable → gate",
         "That stored variable is the one tested in the\nbranch instruction guarding the sink call."),
        ("3  gate → sink",
         "When the branch is satisfied, execution reaches\nthe sink with no additional blocking condition."),
        ("4  input → argument",
         "The HTTP-supplied payload becomes the\n(untransformed) argument to the sink function."),
    ]
    cols = [ACCENT, ACCENT2, RGBColor(0xFF,0xCC,0x55), WARN]
    for i, (title, body) in enumerate(links):
        x = 0.45 + i*3.2
        bx = box(sl, x, 1.25, 3.0, 2.8)
        # colour top strip
        strip = sl.shapes.add_shape(1, Inches(x), Inches(1.25), Inches(3.0), Inches(0.22))
        strip.fill.solid(); strip.fill.fore_color.rgb = cols[i]
        strip.line.fill.background()
        txt(sl, title, x+0.1, 1.32, 2.8, 0.38, size=13, bold=True, color=cols[i])
        txt(sl, body,  x+0.1, 1.75, 2.8, 1.8, size=14, color=LIGHT)

    # flow diagram
    nodes = ["Config\noption", "Gate\nvariable", "Branch\nguard", "Sink\ncall", "Argument\n(payload)"]
    nw, nh = 1.7, 0.75
    ny = 4.45
    for k, n in enumerate(nodes):
        nx = 0.6 + k*2.4
        c  = cols[k-1] if 0 < k < 5 else ACCENT
        bx = box(sl, nx, ny, nw, nh, fill=BOX_BG, border=c, border_pt=1.5)
        txt(sl, n, nx+0.05, ny+0.08, nw-0.1, nh-0.12, size=13, color=c, align=PP_ALIGN.CENTER, bold=True)
        if k < 4:
            txt(sl, "──›", nx+nw+0.05, ny+0.22, 0.65, 0.35, size=15, color=DIM, bold=True)

    txt(sl, "Each link eliminates a different class of false positive — a 'dead code' sink, a hardcoded constant, wrong entry point, or wrong config source.",
        0.5, 5.5, 12.3, 0.6, size=14, color=DIM, italic=True)


def slide_taxonomy(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Corpus Taxonomy  (C / I / S axes)",
                "Each sample labelled on three independent axes")

    axes = [
        ("C — Configuration complexity", ACCENT, [
            "C1   — one CLI flag",
            "C1b — two CLI flags",
            "C1c — three or more CLI flags",
            "C2   — config-file key(s)",
            "C3   — config file + runtime HTTP init",
        ]),
        ("I — Input hop count", ACCENT2, [
            "I1  — input received and passed to sink\n       in the same function",
            "I2  — one intermediate function",
            "I3  — two or more intermediate functions\n       (multi-hop taint path)",
        ]),
        ("S — Sanitization", WARN, [
            "S1  — no sanitization present",
            "S2  — sanitization present but bypassable",
            "S3  — config-gated sanitization\n       (must be absent for exploitation)",
        ]),
    ]

    for i, (header, col, items) in enumerate(axes):
        x = 0.45 + i*4.25
        bx = box(sl, x, 1.2, 4.0, 5.5)
        strip = sl.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(4.0), Inches(0.25))
        strip.fill.solid(); strip.fill.fore_color.rgb = col
        strip.line.fill.background()
        txt(sl, header, x+0.12, 1.24, 3.75, 0.38, size=13, bold=True, color=col)
        cy = 1.65
        for item in items:
            txt(sl, item, x+0.18, cy, 3.7, 0.55, size=14, color=LIGHT)
            cy += 0.58

    txt(sl, "Labels are inferred by the scoring layer from agent output — the agent never sees the taxonomy.",
        0.5, 6.9, 12.3, 0.42, size=13, color=DIM, italic=True)


def slide_tools(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Analysis Toolchain", "Tools available to the LLM agent")

    rows = [
        ("asm_function",         "objdump",    "Disassemble a named function — always tried first; no JVM overhead"),
        ("joern_get_variable_defs", "Joern CPG", "Find all assignments to a variable across function boundaries"),
        ("joern_trace_condition","Joern CPG",  "Find dominating branch conditions at a specific address"),
        ("joern_get_callers",    "Joern CPG",  "List all callers of a function (call graph traversal)"),
        ("joern_forward_taint",  "Joern CPG",  "Forward dataflow taint from an HTTP source to the sink"),
        ("binary_strings",       "strings(1)", "Extract printable strings — used to verify flag / key names"),
        ("joern_raw",            "Joern CPG",  "Arbitrary Scala CPG query — last resort only"),
        ("angr_solve_value",     "angr",       "Symbolic execution to find a non-obvious comparison value"),
    ]

    hx = [0.45, 3.2, 5.1]
    for j, h in enumerate(["Tool name", "Backend", "Purpose"]):
        txt(sl, h, hx[j], 1.18, 2.5, 0.32, size=13, bold=True, color=ACCENT)

    strip = sl.shapes.add_shape(1, Inches(0.45), Inches(1.5), Inches(12.4), Inches(0.04))
    strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    for i, (tool, backend, desc) in enumerate(rows):
        y = 1.6 + i*0.6
        fill = BOX_BG if i % 2 == 0 else RGBColor(0x13,0x17,0x2A)
        box(sl, 0.45, y, 12.4, 0.58, fill=fill, border=BOX_BDR, border_pt=0.4)
        txt(sl, tool,    0.55, y+0.1, 2.55, 0.42, size=13, bold=True,  color=ACCENT2)
        txt(sl, backend, 3.25, y+0.1, 1.75, 0.42, size=13, color=DIM)
        txt(sl, desc,    5.15, y+0.1, 7.6,  0.42, size=13, color=LIGHT)


def slide_results(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Evaluation Results  (50 samples)", "After corpus ground-truth fixes and scoring improvements")

    # Big numbers
    metrics = [
        ("Full accuracy",      "91.3%",  ACCENT2),
        ("Verdict accuracy",   "95.7%",  ACCENT),
        ("Prec. recall",       "0.94",   ACCENT),
        ("Prec. precision",    "0.89",   ACCENT),
    ]
    for i, (label, val, col) in enumerate(metrics):
        x = 0.45 + i*3.2
        colored_box_txt(sl, label, val, x, 1.25, 3.0, h=1.45, lcolor=col, vsize=40)

    # C-axis breakdown table
    txt(sl, "Full accuracy by C-axis", 0.45, 3.0, 5.5, 0.38, size=15, bold=True, color=ACCENT)

    rows_c = [
        ("C1",  "10/10", "100%", ACCENT2),
        ("C1b", "7/9",   "77.8%", ACCENT),
        ("C1c", "10/11", "90.9%", ACCENT),
        ("C2",  "6/7",   "85.7%", ACCENT),
        ("C3",  "9/9",   "100%",  ACCENT2),
    ]
    for i, (label, frac, pct, col) in enumerate(rows_c):
        y = 3.45 + i*0.52
        box(sl, 0.45, y, 5.5, 0.48, fill=BOX_BG if i%2==0 else RGBColor(0x13,0x17,0x2A), border=BOX_BDR, border_pt=0.4)
        txt(sl, label, 0.6,  y+0.08, 0.8, 0.35, size=15, bold=True, color=col)
        txt(sl, frac,  1.6,  y+0.08, 1.2, 0.35, size=15, color=LIGHT)
        # bar
        bar_w = float(pct.strip('%')) / 100 * 3.0
        bar = sl.shapes.add_shape(1, Inches(2.8), Inches(y+0.16), Inches(bar_w), Inches(0.22))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        txt(sl, pct, 2.85+bar_w, y+0.08, 0.8, 0.35, size=13, color=col)

    # Right: other axes
    txt(sl, "Other axes", 6.5, 3.0, 5.5, 0.38, size=15, bold=True, color=ACCENT)
    other = [
        ("C-axis accuracy",  "89.1%", ACCENT),
        ("S-axis accuracy",  "52.2%", WARN),
        ("I-axis accuracy",  "43.5%", WARN),
        ("Duration p50",     "374 s", DIM),
        ("Duration p95",     "932 s", DIM),
    ]
    for i, (label, val, col) in enumerate(other):
        y = 3.45 + i*0.52
        box(sl, 6.5, y, 6.3, 0.48, fill=BOX_BG if i%2==0 else RGBColor(0x13,0x17,0x2A), border=BOX_BDR, border_pt=0.4)
        txt(sl, label, 6.65, y+0.08, 4.0, 0.35, size=14, color=LIGHT)
        txt(sl, val,  10.8,  y+0.08, 1.8, 0.35, size=14, bold=True, color=col, align=PP_ALIGN.RIGHT)


def slide_findings(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Key Engineering Findings", "What we learned building and evaluating the system")

    items = [
        ("Corpus ground-truth bug — wrong flag names in 25 samples",
         "Ground truth used template names (--exec-logging, --exec-init) not compiled into the binaries.\n"
         "Fixed by extracting actual flags from PoC scripts. Raised full accuracy +9 points."),
        ("Multi-sink aggregation error",
         "Picking 'fewest preconditions' as canonical selected trivial CGI sinks over the intended exec sink.\n"
         "Fixed in scoring: evaluate all exploitable sinks, credit the best-matching one."),
        ("C3 token matching too strict",
         "Runtime init tokens (e.g. s3cr3t) come from a config file at runtime — not in the binary.\n"
         "Fixed: scoring ignores the token value for runtime_state preconditions; matches on route only."),
        ("Upstream gate tracing gap",
         "Agent only searched for gate conditions in the sink's immediate caller, missing flags checked\n"
         "in dispatcher functions one or two levels up. Fixed via explicit caller-tracing prompt guidance."),
    ]

    for i, (title, body) in enumerate(items):
        y = 1.18 + i*1.5
        bx = box(sl, 0.45, y, 12.4, 1.35)
        strip = sl.shapes.add_shape(1, Inches(0.45), Inches(y), Inches(0.18), Inches(1.35))
        strip.fill.solid()
        strip.fill.fore_color.rgb = [ACCENT, ACCENT2, WARN, RGBColor(0xFF,0xCC,0x55)][i]
        strip.line.fill.background()
        txt(sl, title, 0.78, y+0.1,  11.9, 0.42, size=15, bold=True, color=WHITE)
        txt(sl, body,  0.78, y+0.52, 11.9, 0.72, size=13, color=LIGHT)


def slide_remaining(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Remaining Challenges", "What's still hard and why")

    items = [
        ("S-axis accuracy  52%",       WARN,
         "Distinguishing S1 (no sanitization) from S2 (bypassable) requires the agent's\n"
         "sanitization evidence — now stored in SinkResult; will improve on fresh runs."),
        ("I-axis accuracy  44%",       WARN,
         "hop_count not stored in old cached results; inference only works for fresh runs\n"
         "that include input-path analysis with the updated schema."),
        ("Kore short flags  (-A -B -C)", WARN,
         "Single-char getopt options are not string literals — strings(1) can't find them.\n"
         "Agent must read the getopt() option string in assembly to recover flag chars."),
        ("C1b accuracy  78%",           DIM,
         "Two-flag cases still miss upstream gates in some binaries; caller-tracing prompt\n"
         "improvements should help on fresh runs. Kore samples are structurally harder."),
    ]

    for i, (title, col, body) in enumerate(items):
        y = 1.18 + i*1.48
        bx = box(sl, 0.45, y, 12.4, 1.32)
        strip = sl.shapes.add_shape(1, Inches(0.45), Inches(y), Inches(0.18), Inches(1.32))
        strip.fill.solid(); strip.fill.fore_color.rgb = col
        strip.line.fill.background()
        txt(sl, title, 0.78, y+0.08,  11.9, 0.42, size=15, bold=True, color=col)
        txt(sl, body,  0.78, y+0.50, 11.9, 0.72, size=13, color=LIGHT)


def slide_improvements(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Improvements Made During Evaluation", "Prompt, scoring, and infrastructure changes")

    cols_h = ["Component", "Change", "Impact"]
    cw     = [2.8, 5.8, 3.5]
    cx     = [0.45, 3.35, 9.25]
    for j, (h, w) in enumerate(zip(cols_h, cw)):
        txt(sl, h, cx[j], 1.18, w, 0.32, size=13, bold=True, color=ACCENT)

    strip = sl.shapes.add_shape(1, Inches(0.45), Inches(1.5), Inches(12.4), Inches(0.04))
    strip.fill.solid(); strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    rows = [
        ("trace_and_classify_gate", "Explicit caller-tracing (joern_get_callers + asm on each caller)",          "Finds upstream dispatcher gates"),
        ("trace_and_classify_gate", "Mandatory binary_strings verification before reporting any flag name",       "Eliminates hallucinated flag names"),
        ("causal_verify",           "Sharpened verdict rules: NotExploitable = structural only; wrong values → Inconclusive", "Enables retry on bad hypotheses"),
        ("synthesize_preconditions","Completeness rule for N independent flags; token sourcing from gate_conditions", "Reduces missing-flag rate"),
        ("causal_verify",           "Store sanitization assessment in SinkResult evidence",                        "Enables S2 vs S1 distinction"),
        ("scoring",                 "Best-sink matching: score against all exploitable sinks, take highest recall", "+9 pt full accuracy"),
        ("scoring",                 "Relax runtime_state value matching — token is not statically determinable",   "C3: 0/9 → 9/9"),
        ("state.py",                "_keep_last reducer on binary/joern_cpg fields",                               "Fixes parallel-sink crashes"),
    ]
    for i, (comp, change, impact) in enumerate(rows):
        y = 1.58 + i*0.6
        fill = BOX_BG if i%2==0 else RGBColor(0x13,0x17,0x2A)
        box(sl, 0.45, y, 12.4, 0.57, fill=fill, border=BOX_BDR, border_pt=0.4)
        txt(sl, comp,   cx[0], y+0.1, cw[0]-0.1, 0.4, size=12, bold=True, color=ACCENT2)
        txt(sl, change, cx[1], y+0.1, cw[1]-0.1, 0.4, size=12, color=LIGHT)
        txt(sl, impact, cx[2], y+0.1, cw[2]-0.1, 0.4, size=12, color=DIM)


def slide_infra(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Evaluation Infrastructure", "How we measure the agent at scale")

    left = [
        "eval.py — batch runner with flags:",
        "  --all / -n N / --sample IDs",
        "  --filter C3 / --filter darkhttpd  (AND-able)",
        "  --resume  (skip completed samples)",
        "  --rescore  (re-score cached results, no API calls)",
        "  --workers N  (parallel threads)",
        "  --logs  (per-sample DEBUG log files)",
        "",
        "Result files: eval_results/<sample>.json",
        "  Contains: AnalysisResult, score dict, duration, log path",
        "",
        "Joern CPG cached per binary — rebuild only once",
        "Per-sample timeout (default 600 s)",
        "Graceful SIGTERM handler — cleans up Joern servers",
    ]
    bullet_block(sl, left, l=0.55, t=1.15, w=6.1, size=16, gap=0.36, leader="")

    # Scoring metrics box
    box(sl, 7.0, 1.15, 5.85, 5.9)
    txt(sl, "Scoring metrics", 7.15, 1.22, 5.5, 0.38, size=15, bold=True, color=ACCENT)
    metrics = [
        "Verdict accuracy  (Exploitable / NotExploitable / Inconclusive)",
        "Precondition recall  (fraction of required flags found)",
        "Precondition precision  (fraction of predicted flags that are correct)",
        "Full accuracy  (verdict correct AND recall = 1.0)",
        "C-axis / S-axis / I-axis accuracy",
        "False-exploitable / false-not-exploitable / gave-up counts",
        "Duration p50 / p95",
    ]
    cy = 1.65
    for m in metrics:
        txt(sl, "›  " + m, 7.15, cy, 5.55, 0.42, size=13, color=LIGHT)
        cy += 0.44

    txt(sl, "Ground truth: ground_truth.json per sample\nPoC.sh: verified exploit script (source of truth for flag names)",
        7.15, 6.15, 5.55, 0.7, size=12, color=DIM, italic=True)


def slide_next(prs):
    sl = blank_slide(prs)
    slide_title(sl, "Next Steps", "What to tackle next")

    items = [
        (ACCENT,  "Run full corpus",
         "Fresh 190-sample run with all prompt and scoring improvements to get clean baseline numbers.\n"
         "Especially needed for S-axis and I-axis which rely on data only stored in new runs."),
        (ACCENT2, "Kore getopt awareness",
         "Teach the agent to extract single-char getopt option strings from the binary (e.g. read\n"
         "the getopt() call's option-string argument in assembly) to find -A/-B/-C style flags."),
        (WARN,    "S-axis improvement",
         "With sanitization evidence now stored in results, fresh runs should distinguish S1/S2.\n"
         "Investigate S3→S1 misses where config-gated sanitization isn't detected."),
        (RGBColor(0xFF,0xCC,0x55), "Remaining corpus ground truth",
         "Kore samples still have unverified flag names; lwan/h2o may need C-axis reclassification.\n"
         "Run the PoC scripts directly to confirm exactly which flags each binary requires."),
    ]

    for i, (col, title, body) in enumerate(items):
        y = 1.18 + i*1.48
        bx = box(sl, 0.45, y, 12.4, 1.32)
        strip = sl.shapes.add_shape(1, Inches(0.45), Inches(y), Inches(0.18), Inches(1.32))
        strip.fill.solid(); strip.fill.fore_color.rgb = col
        strip.line.fill.background()
        num = sl.shapes.add_shape(1, Inches(0.72), Inches(y+0.38), Inches(0.38), Inches(0.38))
        num.fill.solid(); num.fill.fore_color.rgb = col
        num.line.fill.background()
        txt(sl, str(i+1), 0.72, y+0.38, 0.38, 0.38, size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
        txt(sl, title, 1.22, y+0.08,  11.4, 0.42, size=15, bold=True, color=col)
        txt(sl, body,  1.22, y+0.50, 11.4, 0.72, size=13, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# Build
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_cover(prs)
    slide_problem(prs)
    slide_approach(prs)
    slide_causal(prs)
    slide_taxonomy(prs)
    slide_tools(prs)
    slide_results(prs)
    slide_findings(prs)
    slide_remaining(prs)
    slide_improvements(prs)
    slide_infra(prs)
    slide_next(prs)

    out = "/home/ubuntu/anthropic/binary_vuln_agent.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
