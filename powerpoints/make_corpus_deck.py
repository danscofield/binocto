"""Two-slide deck: how the evaluation corpus was built using Claude."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (matches main deck) ───────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x0F, 0x1A)
ACCENT  = RGBColor(0x5B, 0x8D, 0xFF)
ACCENT2 = RGBColor(0x7C, 0xE8, 0xB4)
WARN    = RGBColor(0xFF, 0x7C, 0x5C)
GOLD    = RGBColor(0xFF, 0xCC, 0x55)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xCC, 0xD6, 0xF1)
DIM     = RGBColor(0x77, 0x88, 0xAA)
BOX_BG  = RGBColor(0x1A, 0x1E, 0x35)
BOX_BDR = RGBColor(0x2E, 0x3A, 0x5C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Primitives ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = BG

def rect(slide, l, t, w, h, fill=BOX_BG, border=BOX_BDR, bpt=1.0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(bpt)
    return s

def bar(slide, l, t, w, h, color=ACCENT):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def txt(slide, text, l, t, w, h, size=16, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def title_block(slide, title, subtitle=None):
    bg(slide)
    bar(slide, 0.5, 0.58, 12.33, 0.06)
    txt(slide, title,    0.5, 0.1,  12.33, 0.55, size=32, bold=True)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.68, 12.33, 0.38, size=16, color=LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Corpus design
# ═══════════════════════════════════════════════════════════════════════════════

def slide_corpus_design(prs):
    sl = blank(prs)
    title_block(sl,
        "The Evaluation Corpus  —  188 Synthetic Vulnerabilities",
        "Orthogonal 3-axis taxonomy across 15 real open-source HTTP servers")

    # ── Axis boxes (top row) ──────────────────────────────────────────────────
    axes = [
        (ACCENT,  "C — Configuration gate",
         "C1  · Single CLI flag\n"
         "C1b · Two CLI flags\n"
         "C1c · Three CLI flags\n"
         "C2  · Config-file key\n"
         "C3  · Config file + prior HTTP call",
         "How many configuration conditions must an attacker satisfy to\n"
         "open the code path leading to the vulnerable sink?"),
        (ACCENT2, "I — Input indirection",
         "I1  · Direct: cmd passed verbatim to system()\n"
         "I2  · Buffered: cmd wrapped in snprintf shell string\n"
         "I3  · Struct dispatch: cmd carried through struct + helper fn",
         "How many transforms separate the attacker-controlled HTTP\n"
         "parameter from the argument to system()?"),
        (GOLD,    "S — Sanitization",
         "S1  · None (no filter)\n"
         "S2  · Bypassable (blocks '|'; bypass with ';' or '>')\n"
         "S3  · Config-gated strict filter (exploitable only when\n"
         "      --strict-exec is absent)",
         "Is there a sanitization check? Can it be bypassed, and\n"
         "does configuration control whether it runs at all?"),
    ]

    bw = 4.0
    for i, (col, title, levels, question) in enumerate(axes):
        x = 0.42 + i * (bw + 0.24)
        bx = rect(sl, x, 1.22, bw, 3.65)
        bar(sl, x, 1.22, bw, 0.28, color=col)
        txt(sl, title, x+0.14, 1.25, bw-0.22, 0.26, size=12, bold=True, color=BG)
        txt(sl, levels, x+0.14, 1.55, bw-0.22, 1.55, size=12, color=col)
        bar(sl, x+0.14, 3.08, bw-0.28, 0.02, color=BOX_BDR)
        txt(sl, question, x+0.14, 3.14, bw-0.22, 0.65, size=11, color=DIM, italic=True)

    # ── Coverage stat ─────────────────────────────────────────────────────────
    bx = rect(sl, 0.42, 5.0, 12.49, 0.55,
              fill=RGBColor(0x0C, 0x14, 0x28), border=ACCENT, bpt=1.0)
    txt(sl,
        "5 C-levels  \u00d7  3 I-levels  \u00d7  3 S-levels  =  45 axis combinations  "
        "\u2014  all represented across 15 real servers (darkhttpd, civetweb, lighttpd, lwan, h2o, mongoose, kore, onion, ulfius, libwebsockets, seasocks, ...)",
        0.6, 5.06, 12.1, 0.44, size=13, color=ACCENT, bold=True)

    # ── Per-sample artifacts ──────────────────────────────────────────────────
    txt(sl, "Every sample contains:", 0.42, 5.72, 4.2, 0.32, size=13, bold=True, color=LIGHT)

    artifacts = [
        (ACCENT,  "patch.diff",    "Unified diff injecting the synthetic vulnerability into unmodified server source"),
        (ACCENT2, "build.sh",      "Deterministic compile script; reproduces binary from any clean checkout"),
        (GOLD,    "PoC.sh",        "End-to-end exploit: starts server, triggers RCE, verifies marker; negative control included"),
        (WARN,    "PROVENANCE",    "Structured metadata: axes, gate mechanism, sink location, exploit recipe, changed files"),
        (DIM,     "binary",        "Pre-built patched executable — evaluation never requires recompilation"),
    ]

    ax = 0.42; ay = 6.08; aw = 2.38
    for col, name, desc in artifacts:
        bx = rect(sl, ax, ay, aw, 0.8)
        bar(sl, ax, ay, 0.12, 0.8, color=col)
        txt(sl, name, ax+0.22, ay+0.06, aw-0.28, 0.28, size=12, bold=True, color=col)
        txt(sl, desc, ax+0.22, ay+0.36, aw-0.28, 0.36, size=10, color=LIGHT)
        ax += aw + 0.15


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Built by Claude
# ═══════════════════════════════════════════════════════════════════════════════

def slide_built_by_claude(prs):
    sl = blank(prs)
    title_block(sl,
        "Built by Claude  —  21 Parallel Agents, ~65 Minutes",
        "The corpus was itself produced by the same agentic approach it evaluates")

    # ── Left: the challenge ───────────────────────────────────────────────────
    txt(sl, "What corpus construction required", 0.5, 1.12, 5.9, 0.38,
        size=15, bold=True, color=WARN)

    challenges = [
        ("Deep server expertise",
         "Each of 15 HTTP servers has its own request dispatch loop, "
         "argv parsing idioms, config-file format, and build system. "
         "Creating a correct patch requires reading and understanding all of them."),
        ("Correct, verified patches",
         "Every diff must: compile cleanly, route HTTP correctly, "
         "wire the gate variable to the exact right flag/file key, "
         "and produce a system() call reachable through the intended I-axis path."),
        ("End-to-end PoC scripts",
         "Each PoC must start the patched server, send a working exploit payload, "
         "verify RCE via a marker file, then re-run with the gate disabled "
         "and confirm the endpoint returns 404 (negative control)."),
        ("940 files, 45 axis combinations",
         "188 samples x 5 artifacts = 940 files. Every C x I x S combination must "
         "be represented across diverse servers with no copy-paste shortcuts "
         "— each server's code path is unique."),
    ]
    cy = 1.52
    for name, body in challenges:
        bx = rect(sl, 0.5, cy, 5.9, 1.2)
        bar(sl, 0.5, cy, 0.14, 1.2, color=WARN)
        txt(sl, name, 0.78, cy+0.08, 5.5, 0.36, size=13, bold=True, color=WARN)
        txt(sl, body, 0.78, cy+0.46, 5.5, 0.65, size=11.5, color=LIGHT)
        cy += 1.28

    # ── Divider ───────────────────────────────────────────────────────────────
    bar(sl, 6.62, 1.10, 0.04, 5.85, color=BOX_BDR)

    # ── Right: how Claude did it ───────────────────────────────────────────────
    txt(sl, "How Claude built it", 6.82, 1.12, 6.0, 0.38,
        size=15, bold=True, color=ACCENT2)

    waves = [
        (ACCENT,  "Wave 1 — 15 parallel agents",
         "One agent per server. Each agent read the server source tree, "
         "identified the request dispatch loop and argv parser, "
         "wrote a patch for C1/C1b/C1c x I1-I3 x S1-S3 variants, "
         "compiled each binary, and ran the PoC to confirm RCE."),
        (ACCENT2, "Wave 2 — 6 parallel agents",
         "Six agents added C2 and C3 variants across the same servers: "
         "config-file-parsed gates (C2) and runtime-state gates requiring a "
         "prior HTTP call to /exec/init with a secret token (C3). "
         "72 new samples; each verified end-to-end."),
        (GOLD,    "Per-agent workflow",
         "Read source \u2192 patch \u2192 build.sh \u2192 PoC.sh \u2192 compile \u2192 "
         "exploit (RCE confirmed) \u2192 negative control (gate disabled, 404 returned) "
         "\u2192 write PROVENANCE + ground_truth.json."),
        (DIM,     "Total wall time: ~65 minutes",
         "188 samples, 940 files, 15 servers, 45 axis combinations. "
         "A task that would take a team of expert security researchers "
         "several weeks was completed in a single afternoon."),
    ]
    cy2 = 1.52
    cols2 = [ACCENT, ACCENT2, GOLD, DIM]
    for col, name, body in waves:
        bx = rect(sl, 6.82, cy2, 6.0, 1.2)
        bar(sl, 6.82, cy2, 0.14, 1.2, color=col)
        txt(sl, name, 7.1, cy2+0.08, 5.6, 0.36, size=13, bold=True, color=col)
        txt(sl, body, 7.1, cy2+0.46, 5.6, 0.65, size=11.5, color=LIGHT)
        cy2 += 1.28

    # ── Bottom callout ────────────────────────────────────────────────────────
    bx = rect(sl, 0.45, 6.88, 12.4, 0.52,
              fill=RGBColor(0x08, 0x14, 0x10), border=ACCENT2, bpt=1.2)
    txt(sl,
        "The corpus is a demonstration of the approach, not just an evaluation of it.  "
        "Claude read 15 real server codebases, synthesised 188 precisely structured vulnerabilities "
        "across an orthogonal 3-axis design, verified each one end-to-end, and produced "
        "940 machine-readable artifacts — all without human intervention in the construction loop.",
        0.6, 6.92, 12.1, 0.44, size=12.5, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# Build
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_corpus_design(prs)
    slide_built_by_claude(prs)
    out = "/home/ubuntu/anthropic/corpus_built_by_claude.pptx"
    prs.save(out)
    print(f"Saved \u2192 {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
